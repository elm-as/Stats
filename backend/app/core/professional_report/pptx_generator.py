"""
Génération PPTX — style présentation exécutive.
"""

from __future__ import annotations

import io
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.core.professional_report.builder import ReportContent, ReportSection


ACCENT = RGBColor(0x06, 0xB6, 0xD4)
SECONDARY = RGBColor(0x1E, 0x3A, 0x8A)
GRAY_TXT = RGBColor(0x1F, 0x29, 0x37)
GRAY_LIGHT = RGBColor(0x6B, 0x72, 0x80)

SEVERITY_RGB = {
    "critical": RGBColor(0xDC, 0x26, 0x26),
    "warning": RGBColor(0xD9, 0x77, 0x06),
    "success": RGBColor(0x05, 0x96, 0x69),
    "methodological": RGBColor(0x7C, 0x3A, 0xED),
    "info": RGBColor(0x25, 0x63, 0xEB),
}


def _strip_md(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text


def _add_textbox(slide, x, y, w, h, text, *,
                 size=14, bold=False, color=GRAY_TXT, align=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = _strip_md(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return box


def _add_bar(slide, x, y, w, h, color):
    """Trait coloré (header band)."""
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # rect
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _slide_cover(prs: Presentation, content: ReportContent):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Bande haute
    _add_bar(slide, Emu(0), Emu(0), prs.slide_width, Inches(0.6), ACCENT)

    _add_textbox(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
                 content.title, size=40, bold=True, color=SECONDARY, align=PP_ALIGN.CENTER,
                 font_name="Calibri Light")

    _add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.7),
                 content.subtitle, size=20, color=ACCENT, align=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
                 f"{content.date}  •  {content.author}", size=11, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)


def _slide_section_title(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bar(slide, Emu(0), Emu(0), prs.slide_width, Inches(0.4), ACCENT)
    _add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(2),
                 title, size=44, bold=True, color=SECONDARY, align=PP_ALIGN.CENTER,
                 font_name="Calibri Light")


def _slide_bullets(prs: Presentation, title: str, bullets: list[str], subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bar(slide, Emu(0), Emu(0), prs.slide_width, Inches(0.3), ACCENT)

    _add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
                 title, size=24, bold=True, color=SECONDARY)

    if subtitle:
        _add_textbox(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
                     subtitle, size=12, color=GRAY_LIGHT)

    # Bullets dans une seule textbox
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12.0), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        r = p.add_run()
        r.text = f"●  {_strip_md(b)}"
        r.font.size = Pt(14)
        r.font.color.rgb = GRAY_TXT
        r.font.name = "Calibri"
        p.space_after = Pt(8)


def _slide_insights(prs: Presentation, title: str, insights: list[dict], max_per_slide: int = 6):
    """Découpe les insights en plusieurs slides si besoin."""
    chunks = [insights[i:i + max_per_slide] for i in range(0, len(insights), max_per_slide)]
    for ci, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bar(slide, Emu(0), Emu(0), prs.slide_width, Inches(0.3), ACCENT)

        suffix = f" ({ci + 1}/{len(chunks)})" if len(chunks) > 1 else ""
        _add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
                     title + suffix, size=24, bold=True, color=SECONDARY)

        # Affichage en grille 2 colonnes
        for idx, ins in enumerate(chunk):
            col = idx % 2
            row = idx // 2
            x = Inches(0.5 + col * 6.3)
            y = Inches(1.4 + row * 1.8)

            sev = ins.get("severity", "info")
            sev_col = SEVERITY_RGB.get(sev, SEVERITY_RGB["info"])

            # Cartouche bordure
            shape = slide.shapes.add_shape(5, x, y, Inches(6.1), Inches(1.6))  # rounded rect
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
            shape.line.color.rgb = sev_col
            shape.line.width = Pt(1.5)

            # Texte intérieur
            box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                            Inches(5.8), Inches(1.4))
            tf = box.text_frame
            tf.word_wrap = True

            p1 = tf.paragraphs[0]
            r1 = p1.add_run()
            r1.text = _strip_md(ins.get("title", ""))
            r1.font.bold = True
            r1.font.size = Pt(12)
            r1.font.color.rgb = sev_col

            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = _strip_md(ins.get("message", ""))[:200]
            r2.font.size = Pt(9.5)
            r2.font.color.rgb = GRAY_TXT


def _slide_table(prs: Presentation, title: str, headers: list, rows: list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bar(slide, Emu(0), Emu(0), prs.slide_width, Inches(0.3), ACCENT)
    _add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.6),
                 title, size=24, bold=True, color=SECONDARY)

    n_rows = min(len(rows), 12) + 1
    n_cols = len(headers)
    if n_rows < 2 or n_cols < 1:
        return

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.5),
        Inches(12.3), Inches(min(5.5, 0.5 * n_rows)),
    ).table

    for ci, h in enumerate(headers):
        cell = table_shape.cell(0, ci)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(h)
        r.font.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for ri, row in enumerate(rows[:12]):
        for ci, val in enumerate(row[:n_cols]):
            cell = table_shape.cell(ri + 1, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(10)
            r.font.color.rgb = GRAY_TXT


def generate_pptx(content: ReportContent) -> bytes:
    """Génère un PPTX depuis un ReportContent."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cover
    _slide_cover(prs, content)

    # Résumé exécutif
    if content.executive_summary or content.key_findings:
        bullets = []
        if content.executive_summary:
            bullets.append(content.executive_summary)
        bullets.extend(content.key_findings or [])
        _slide_bullets(prs, "Résumé exécutif", bullets, subtitle="Vue d'ensemble")

    # Sections
    for sec in content.sections:
        # Insights en grille
        if sec.insights:
            _slide_insights(prs, sec.title, sec.insights)
        elif sec.bullets:
            _slide_bullets(prs, sec.title, sec.bullets, subtitle=sec.body or "")
        elif sec.table:
            _slide_table(prs, sec.title, sec.table.get("headers", []), sec.table.get("rows", []))
        elif sec.body:
            _slide_bullets(prs, sec.title, [sec.body])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
